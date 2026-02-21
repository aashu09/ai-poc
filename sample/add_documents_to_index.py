import os
import io
import fitz  # PyMuPDF for PDFs
import pandas as pd
from keybert import KeyBERT
from langdetect import detect
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
import requests
import json
import numpy as np
import time
from datetime import datetime
from azure.storage.filedatalake import (
    DataLakeServiceClient
)
from azure.identity import ClientSecretCredential
from langchain.text_splitter import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
from openai import AzureOpenAI
from pathlib import Path
from core.config import Settings

env_path = Path('.') / '.env'
load_dotenv(dotenv_path=env_path)

endpoint = Settings.openai_endpoint
if not endpoint:
    endpoint = os.getenv("openai_endpoint")

model_name = Settings.model_name
if not model_name:
    model_name = os.getenv("model_name")

deployment = Settings.deployment
if not deployment:
    deployment = os.getenv("deployment")

subscription_key = Settings.subscription_key
if not subscription_key:
    subscription_key = os.getenv("subscription_key")

api_version = Settings.api_version
if not api_version:
    api_version = os.getenv("api_version")

client = AzureOpenAI(
    api_version=api_version,
    azure_endpoint=endpoint,
    api_key=subscription_key,
)
# === CONFIG ===

# search_service = "<your-search-service-name>"
# search_index = "lakehouse-index"
# admin_key = "<your-admin-key>"
# embedding_api_url = "<your-embedding-api-endpoint>"

search_service = Settings.search_service_name
if not search_service:
    search_service = os.getenv("search_service")

search_index = "index-graphpoc-001"

search_service_admin_key = Settings.search_service_admin_key
if not search_service_admin_key:
    search_service_admin_key = os.getenv("search_service_admin_key")

embedding_api_url = Settings.embedding_model_url
if not embedding_api_url:
    embedding_api_url = os.getenv("embedding_model_url")

embedding_model_key = Settings.embedding_model_key
if not embedding_model_key:
    embedding_model_key = os.getenv("embedding_model_key")

# Set your account, workspace, and item path here
ACCOUNT_NAME = "onelake"
WORKSPACE_NAME = "workspace-aifoundary-poc"
DATA_PATH = "aifoundary_lakehouse_v1.Lakehouse/Files/raw/data"

# Azure Cognitive Search client
search_client = SearchClient(
    endpoint=f"https://{search_service}.search.windows.net",
    index_name=search_index,
    credential=AzureKeyCredential(search_service_admin_key)
)

account_url = f"https://{ACCOUNT_NAME}.dfs.fabric.microsoft.com"
# token_credential = DefaultAzureCredential()
token_credential = ClientSecretCredential(
    tenant_id="<id>",
    client_id="<client_id>",
    client_secret="<secrete>"
)

service_client = DataLakeServiceClient(account_url, credential=token_credential)

# Create a file system client for the workspace
file_system_client = service_client.get_file_system_client(WORKSPACE_NAME)

# List a directory within the filesystem
paths = file_system_client.get_paths(path=DATA_PATH)

# Initialize KeyBERT
# kw_model = KeyBERT()

# === FUNCTIONS ===

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=24000,
    chunk_overlap=200,
    separators=["\n\n", "\n", ".", " ", ""]
)


def extract_text_from_pdf(file_path):
    file_bytes = file_system_client.get_file_client(file_path).download_file().readall()
    doc = fitz.open("pdf", file_bytes)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_text_from_excel(file_path):
    text = ""
    local_file = io.BytesIO(file_system_client.get_file_client(file_path).download_file().readall())

    try:
        if file_path.lower().endswith(".xls"):
            xls = pd.ExcelFile(local_file, engine="xlrd")
        else:
            xls = pd.ExcelFile(local_file)
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet)
            text += df.to_string(index=False, header=False) + "\n"
    except Exception as e:
        print(f"Error reading Excel file {file_path}: {e}")
    return text


def extract_text_from_csv(file_path):
    text = ""
    local_file = io.BytesIO(file_system_client.get_file_client(file_path).download_file().readall())
    try:
        df = pd.read_csv(local_file)
        text += df.to_string(index=False, header=False) + "\n"
    except Exception as e:
        print(f"❌ Error reading CSV file {file_path}: {e}")
    return text


def extract_text_from_word(file_path):
    try:
        import docx
        word_bytes = file_system_client.get_file_client(file_path).download_file().readall()
        doc = docx.Document(io.BytesIO(word_bytes))
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading Word file {file_path}: {e}")
        return ""


def extract_text_from_ppt(file_path):
    try:
        from pptx import Presentation
        ppt_bytes = file_system_client.get_file_client(file_path).download_file().readall()
        prs = Presentation(io.BytesIO(ppt_bytes))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        print(f"Error reading PPT file {file_path}: {e}")
        return ""


def split_text_into_chunks(text, chunk_size=24000):
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


def get_embedding(text):
    payload = {"input": text}

    headers = {
        "Authorization": f"Bearer {embedding_model_key}",
        "Content-Type": "application/json"
    }
    response = requests.post(embedding_api_url, headers=headers, data=json.dumps(payload))
    if response.status_code == 200:
        return response.json()["data"][0].get("embedding")
    elif response.status_code == 429:
        retry_after = int(response.headers.get('Retry-After', 10))
        time.sleep(retry_after)
        return get_embedding(text)
    else:
        print(f"Embedding API Error: {response.text}")
        return None


def extract_key_phrases_with_openai(text):
    """Extract key phrases using Azure AI Foundry model."""
    try:
        # Define the prompt to extract key phrases
        prompt = f"Extract key phrases from the following text:\n\n{text} \n\n Note: don't include answer like 'Here are the key phrases extracted from the provided text: \n\n Only provide key phrases comma seprated"

        # Send a completion request to Azure OpenAI (your deployed model)
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a helpful assistant.",
                },
                {
                    "role": "user",
                    "content": prompt,
                }
            ],
            max_tokens=200,
            temperature=0.7,
            top_p=1.0,
            model=deployment
        )

        # Parse the response and extract key phrases
        key_phrases = response.choices[0].message.content

        # Convert the key phrases into a list (splitting by commas)
        key_phrases_list = [phrase.strip() for phrase in key_phrases.split(",") if phrase.strip()]

        return key_phrases_list
    except Exception as e:
        print(f"Error extracting key phrases: {e}")
        return []


def detect_language(text):
    try:
        return detect(text)
    except Exception as e:
        print(f" Language detection failed: {e}")
        return "unknown"


def process_file(file_path, file_name):
    file_extension = os.path.splitext(file_name)[-1].lower()
    content_text = ""

    if file_extension == ".pdf":
        content_text = extract_text_from_pdf(file_path)
    elif file_extension in [".xls", ".xlsx"]:
        content_text = extract_text_from_excel(file_path)
    elif file_extension == ".csv":
        content_text = extract_text_from_csv(file_path)
    elif file_extension == ".docx":
        content_text = extract_text_from_word(file_path)
    elif file_extension == ".pptx":
        content_text = extract_text_from_ppt(file_path)
    else:
        print(f" Unsupported file type: {file_extension}")
        return []

    if not content_text.strip():
        print(f"Empty content in {file_name}, skipping.")
        return []

    text_chunks = text_splitter.split_text(content_text)

    file_properties_client = file_system_client.get_file_client(file_path)
    file_properties = file_properties_client.get_file_properties()
    last_modified = file_properties["last_modified"].strftime("%Y-%m-%dT%H:%M:%SZ")

    metadata_storage_path = file_path  # or a URL if uploading later

    documents = []
    for idx, chunk in enumerate(text_chunks):
        language = detect_language(chunk)
        print(f"language:{language}")

        embedding = get_embedding(chunk)

        if not embedding:
            continue

        key_phrases = extract_key_phrases_with_openai(chunk)
        print(f"key_phrases:{key_phrases}")

        file_name_without_extention, _ = os.path.splitext(file_name)

        doc_id = f"{file_name_without_extention}-{idx}" if len(text_chunks) > 1 else file_name

        doc = {
            "id": doc_id,
            "fileName": file_name,
            "fileType": file_extension.replace(".", ""),
            "metadata_storage_path": metadata_storage_path,
            "lastModified": last_modified,
            "language": language,
            "keyPhrases": key_phrases,
            "embedding": embedding,
            "content": chunk
        }
        documents.append(doc)

    return documents


# === MAIN ===

batch = []
batch_size = 2

for path in paths:
    file_name = os.path.basename(path.name)
    file_path = path.name
    print(f"Processing: {file_name}")
    print(f"file_path: {file_path}")

    try:
        docs = process_file(file_path, file_name)
        if docs:
            batch.extend(docs)

            if len(batch) >= batch_size:
                search_client.merge_or_upload_documents(documents=batch)
                print(f"Uploaded batch of {len(batch)} documents")
                batch = []
                time.sleep(1)

    except Exception as e:
        print(f"Error processing {file_name}: {e}")

if batch:
    search_client.merge_or_upload_documents(documents=batch)
    print(f"Uploaded final batch of {len(batch)} documents")

print("All done!")